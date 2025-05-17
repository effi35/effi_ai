#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
מודול ניתוח מסמכים למאחד קוד חכם Pro 2.0
מספק יכולות מתקדמות לניתוח מסמכים טקסטואליים, קבצי מדיה ומסמכי קוד

מחבר: Claude AI
גרסה: 1.0.0
תאריך: מאי 2025
"""

import os
import re
import sys
import json
import base64
import hashlib
import logging
import tempfile
import mimetypes
import collections
import concurrent.futures
from pathlib import Path
from typing import Dict, List, Tuple, Any, Optional, Union, Set, Callable

try:
    # ספריות לעיבוד טקסט ולמידת מכונה
    import numpy as np
    import pandas as pd
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.cluster import KMeans, DBSCAN
    from sklearn.metrics.pairwise import cosine_similarity
    import nltk
    from nltk.tokenize import word_tokenize, sent_tokenize
    from nltk.corpus import stopwords
    from nltk.stem import WordNetLemmatizer
    from gensim.summarization import summarize, keywords
    NLP_AVAILABLE = True
    
    # וידוא קיום של חבילות NLTK הדרושות
    try:
        nltk.data.find('tokenizers/punkt')
    except LookupError:
        nltk.download('punkt', quiet=True)
    try:
        nltk.data.find('corpora/stopwords')
    except LookupError:
        nltk.download('stopwords', quiet=True)
    try:
        nltk.data.find('corpora/wordnet')
    except LookupError:
        nltk.download('wordnet', quiet=True)
        
except ImportError:
    NLP_AVAILABLE = False

try:
    # ספריות לעיבוד תמונה
    from PIL import Image
    import pytesseract
    IMAGE_PROCESSING_AVAILABLE = True
except ImportError:
    IMAGE_PROCESSING_AVAILABLE = False

try:
    # ספריות להמרת PDF וקבצי Office
    import fitz  # PyMuPDF
    import docx
    import openpyxl
    import pptx
    DOCUMENT_CONVERSION_AVAILABLE = True
except ImportError:
    DOCUMENT_CONVERSION_AVAILABLE = False

try:
    # ספריית עיבוד שפה טבעית עברית
    import hazm
    HEBREW_NLP_AVAILABLE = True
except ImportError:
    HEBREW_NLP_AVAILABLE = False

# הגדרת לוגים
logger = logging.getLogger(__name__)

class DocumentAnalyzer:
    """
    מנתח מסמכים חכם - מספק יכולות ניתוח, סיכום, חילוץ מידע ועיבוד מסמכים
    
    יכולות:
    - ניתוח וסיכום תוכן טקסטואלי
    - חילוץ מילות מפתח, נושאים ומושגים מרכזיים
    - חילוץ טקסט מתמונות וקבצי PDF
    - ניתוח דמיון וקשר בין מסמכים
    - המרה בין פורמטים שונים
    - זיהוי שפה וניתוח מסמכים בעברית
    - חילוץ מידע מובנה ומטא-דאטה ממסמכים
    """
    
    def __init__(self, config: dict = None):
        """אתחול מנתח המסמכים עם הגדרות אופציונליות"""
        self.config = config or {}
        self.max_workers = self.config.get("max_workers", 4)
        self.summary_ratio = self.config.get("summary_ratio", 0.2)
        self.keywords_ratio = self.config.get("keywords_ratio", 0.1)
        self.min_similarity = self.config.get("min_similarity", 0.7)
        self.default_language = self.config.get("default_language", "hebrew")
        self.output_dir = self.config.get("output_dir", "output")
        self.cache_dir = self.config.get("cache_dir", "cache")
        
        # יצירת תיקיות נדרשות
        os.makedirs(self.output_dir, exist_ok=True)
        os.makedirs(self.cache_dir, exist_ok=True)
        
        # ודא מיפויי MIME תקינים
        self._ensure_mime_types()
        
        # אתחול רכיבי NLP
        self.vectorizer = None
        self.lemmatizer = None
        self.stop_words = None
        
        if NLP_AVAILABLE:
            self._init_nlp_components()
        
        # מיפוי פונקציות עיבוד לפי סוגי קבצים
        self.processors = self._init_file_processors()
        
        logger.info(f"מנתח מסמכים אותחל: NLP זמין={NLP_AVAILABLE}, "
                   f"עיבוד תמונה זמין={IMAGE_PROCESSING_AVAILABLE}, "
                   f"המרת מסמכים זמינה={DOCUMENT_CONVERSION_AVAILABLE}")

    def _ensure_mime_types(self) -> None:
        """וידוא שכל מיפויי MIME הנדרשים רשומים"""
        # הוספת מיפויים חסרים
        if not mimetypes.inited:
            mimetypes.init()
        
        # הוספת מיפויי קבצי Office אם חסרים
        office_types = {
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.md': 'text/markdown',
            '.json': 'application/json',
            '.yml': 'application/x-yaml',
            '.yaml': 'application/x-yaml'
        }
        
        for ext, mime_type in office_types.items():
            if not mimetypes.guess_type(f"test{ext}")[0]:
                mimetypes.add_type(mime_type, ext)

    def _init_nlp_components(self) -> None:
        """אתחול רכיבי עיבוד שפה טבעית"""
        try:
            # וקטוריזציה של טקסט
            self.vectorizer = TfidfVectorizer(
                max_features=1000,
                stop_words=None,  # נגדיר ידנית לפי שפה
                ngram_range=(1, 2),
                use_idf=True
            )
            
            # עיבוד מורפולוגי
            self.lemmatizer = WordNetLemmatizer()
            
            # הגדרת מילות עצירה לפי שפה
            if self.default_language == "hebrew" and HEBREW_NLP_AVAILABLE:
                self.stop_words = set(hazm.stopwords_list())
            else:
                self.stop_words = set(stopwords.words('english'))
                # הוספת מילות עצירה בעברית
                hebrew_stop_words = self._get_hebrew_stopwords()
                self.stop_words.update(hebrew_stop_words)
        
        except Exception as e:
            logger.error(f"שגיאה באתחול רכיבי NLP: {str(e)}")
            self.vectorizer = None
            self.lemmatizer = None
            self.stop_words = set()

    def _get_hebrew_stopwords(self) -> Set[str]:
        """החזרת רשימת מילות עצירה בעברית"""
        # רשימה בסיסית של מילות עצירה בעברית
        hebrew_stop_words = {
            "אני", "את", "אתה", "הוא", "היא", "אנחנו", "אתם", "הם", "הן",
            "של", "את", "על", "עם", "זה", "זאת", "זו", "אלה", "אלו", 
            "או", "גם", "כי", "אם", "אבל", "לא", "רק", "מן", "אל", "כמו",
            "כך", "לכן", "בגלל", "לפי", "כש", "כאשר", "ש", "ב", "ל", "מ", 
            "יש", "אין", "היה", "יהיה", "להיות", "כל", "כן", "לך", "לי",
            "לנו", "להם", "אותו", "אותה", "אותם", "אותן", "אותי", "אותך",
            "אותנו", "אז", "עכשיו", "פה", "שם", "איך", "מתי", "למה", "מי",
            "מה", "איפה", "מדוע", "כמה", "איזה", "איזו", "אחרי", "לפני",
            "תחת", "מעל", "מתחת", "ליד", "בין", "בתוך", "בלי", "עד", "מאז"
        }
        return hebrew_stop_words

    def _init_file_processors(self) -> Dict[str, Callable]:
        """אתחול מיפוי בין סוגי קבצים לפונקציות עיבוד"""
        processors = {
            # טקסט פשוט
            'text/plain': self._process_text_file,
            'text/markdown': self._process_markdown_file,
            'text/csv': self._process_csv_file,
            'application/json': self._process_json_file,
            'application/x-yaml': self._process_yaml_file,
            
            # מסמכים
            'application/pdf': self._process_pdf_file,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_docx_file,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._process_xlsx_file,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_pptx_file,
            
            # קוד
            'text/x-python': self._process_python_file,
            'application/javascript': self._process_javascript_file,
            'text/html': self._process_html_file,
            'text/css': self._process_css_file,
            'text/xml': self._process_xml_file,
            
            # תמונות
            'image/jpeg': self._process_image_file,
            'image/png': self._process_image_file,
            'image/gif': self._process_image_file,
            'image/svg+xml': self._process_svg_file
        }
        return processors

    def analyze_document(self, file_path: str) -> Dict[str, Any]:
        """
        ניתוח מסמך יחיד
        
        Args:
            file_path: נתיב לקובץ לניתוח
            
        Returns:
            מילון עם תוצאות הניתוח
        """
        if not os.path.exists(file_path):
            logger.error(f"קובץ {file_path} לא נמצא")
            return {"status": "error", "error": f"קובץ {file_path} לא נמצא"}
        
        try:
            # זיהוי סוג הקובץ
            mime_type = self._detect_mime_type(file_path)
            logger.info(f"מנתח מסמך: {file_path} (סוג: {mime_type})")
            
            # בחירת מעבד מתאים
            processor = self.processors.get(mime_type)
            if not processor:
                logger.warning(f"אין מעבד מתאים לסוג קובץ {mime_type}, משתמש במעבד כללי")
                processor = self._process_generic_file
            
            # עיבוד הקובץ
            processing_start_time = time.time()
            results = processor(file_path)
            processing_duration = time.time() - processing_start_time
            
            # הוספת מידע כללי
            file_info = self._get_file_info(file_path)
            results.update(file_info)
            results["processing_time"] = processing_duration
            results["mime_type"] = mime_type
            results["status"] = "success"
            
            # ניתוח תוכן טקסטואלי (אם קיים)
            if "text_content" in results and results["text_content"] and NLP_AVAILABLE:
                text_analysis = self._analyze_text(results["text_content"])
                results.update(text_analysis)
            
            return results
            
        except Exception as e:
            logger.error(f"שגיאה בניתוח מסמך {file_path}: {str(e)}")
            return {"status": "error", "error": str(e), "file_path": file_path}

    def analyze_documents(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """
        ניתוח מספר מסמכים במקביל
        
        Args:
            file_paths: רשימת נתיבים לקבצים לניתוח
            
        Returns:
            רשימת תוצאות ניתוח
        """
        if not file_paths:
            logger.warning("לא סופקו קבצים לניתוח")
            return []
        
        logger.info(f"מנתח {len(file_paths)} קבצים במקביל")
        
        # ניתוח קבצים במקביל
        with concurrent.futures.ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            results = list(executor.map(self.analyze_document, file_paths))
        
        # ניתוח קשרים בין מסמכים
        if len(results) > 1 and NLP_AVAILABLE:
            try:
                # איסוף תוכן טקסטואלי מכל המסמכים
                documents_content = []
                for result in results:
                    if result.get("status") == "success" and "text_content" in result:
                        documents_content.append(result["text_content"])
                    else:
                        documents_content.append("")
                
                # ניתוח דמיון בין מסמכים
                similarity_matrix = self._calculate_document_similarity(documents_content)
                
                # הוספת מידע על דמיון ומסמכים קשורים
                for i, result in enumerate(results):
                    if result.get("status") == "success":
                        related_docs = []
                        for j, sim in enumerate(similarity_matrix[i]):
                            if i != j and sim >= self.min_similarity:
                                related_docs.append({
                                    "file_path": file_paths[j],
                                    "similarity": float(sim)
                                })
                        result["related_documents"] = sorted(related_docs, key=lambda x: x["similarity"], reverse=True)
            except Exception as e:
                logger.error(f"שגיאה בניתוח דמיון בין מסמכים: {str(e)}")
        
        return results

    def _detect_mime_type(self, file_path: str) -> str:
        """
        זיהוי סוג MIME של קובץ
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            סוג MIME
        """
        mime_type, _ = mimetypes.guess_type(file_path)
        
        # אם לא זוהה MIME, ננסה לזהות לפי סיומת
        if not mime_type:
            ext = os.path.splitext(file_path)[1].lower()
            
            # שיוך סיומות נפוצות
            ext_to_mime = {
                '.txt': 'text/plain',
                '.md': 'text/markdown',
                '.py': 'text/x-python',
                '.js': 'application/javascript',
                '.html': 'text/html',
                '.css': 'text/css',
                '.json': 'application/json',
                '.xml': 'text/xml',
                '.csv': 'text/csv',
                '.yml': 'application/x-yaml',
                '.yaml': 'application/x-yaml'
            }
            
            mime_type = ext_to_mime.get(ext, 'application/octet-stream')
        
        # ניסיון לזהות קבצי טקסט
        if mime_type == 'application/octet-stream':
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    # קריאת 1024 בתים ראשונים
                    sample = f.read(1024)
                    if not sample:  # קובץ ריק
                        mime_type = 'text/plain'
                    elif not '\0' in sample:  # לא נראה בינארי
                        mime_type = 'text/plain'
            except:
                pass
        
        return mime_type

    def _get_file_info(self, file_path: str) -> Dict[str, Any]:
        """
        קבלת מידע בסיסי על קובץ
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            מילון עם מידע על הקובץ
        """
        try:
            # מידע בסיסי
            stat_info = os.stat(file_path)
            
            # חישוב חתימת קובץ
            file_hash = ""
            with open(file_path, 'rb') as f:
                file_hash = hashlib.md5(f.read()).hexdigest()
            
            return {
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "file_size": stat_info.st_size,
                "file_modified": stat_info.st_mtime,
                "file_created": stat_info.st_ctime,
                "file_hash": file_hash
            }
        except Exception as e:
            logger.error(f"שגיאה בקבלת מידע על קובץ {file_path}: {str(e)}")
            return {"file_path": file_path}

    def _analyze_text(self, text: str) -> Dict[str, Any]:
        """
        ניתוח תוכן טקסטואלי
        
        Args:
            text: טקסט לניתוח
            
        Returns:
            מילון עם תוצאות הניתוח
        """
        if not text or not NLP_AVAILABLE:
            return {}
        
        try:
            # סטטיסטיקה בסיסית
            word_count = len(text.split())
            char_count = len(text)
            
            # חלוקה למשפטים
            sentences = sent_tokenize(text)
            sent_count = len(sentences)
            
            # זיהוי שפה
            language = self._detect_language(text)
            
            # סיכום ומילות מפתח
            summary = ""
            keywords_list = []
            
            try:
                # סיכום טקסט
                if word_count > 100:  # רק אם יש מספיק טקסט
                    summary = summarize(text, ratio=self.summary_ratio)
                
                # מילות מפתח
                keywords_text = keywords(text, ratio=self.keywords_ratio)
                keywords_list = keywords_text.split('\n') if keywords_text else []
            except Exception as e:
                logger.debug(f"שגיאה בסיכום טקסט או חילוץ מילות מפתח: {str(e)}")
            
            # זיהוי ישויות בטקסט (אם NLTK מתקדם זמין)
            entities = self._extract_entities(text, language)
            
            # כריית נושאים
            topics = self._extract_topics(text, language)
            
            # בניית תוצאות
            result = {
                "word_count": word_count,
                "char_count": char_count,
                "sentence_count": sent_count,
                "language": language,
                "summary": summary,
                "keywords": keywords_list,
                "entities": entities,
                "topics": topics
            }
            
            return result
            
        except Exception as e:
            logger.error(f"שגיאה בניתוח טקסט: {str(e)}")
            return {}

    def _detect_language(self, text: str) -> str:
        """
        זיהוי שפת טקסט
        
        Args:
            text: טקסט לזיהוי
            
        Returns:
            קוד שפה
        """
        # זיהוי פשוט לפי נוכחות תווים בעברית
        hebrew_chars = re.findall(r'[\u0590-\u05FF]', text)
        english_chars = re.findall(r'[a-zA-Z]', text)
        
        if len(hebrew_chars) > len(english_chars):
            return "hebrew"
        else:
            return "english"

    def _extract_entities(self, text: str, language: str) -> Dict[str, List[str]]:
        """
        חילוץ ישויות מטקסט
        
        Args:
            text: טקסט לניתוח
            language: שפת הטקסט
            
        Returns:
            מילון עם ישויות לפי סוג
        """
        entities = {
            "persons": [],
            "organizations": [],
            "locations": [],
            "dates": [],
            "emails": [],
            "urls": []
        }
        
        try:
            # חילוץ כתובות אימייל
            emails = re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', text)
            entities["emails"] = list(set(emails))
            
            # חילוץ כתובות URL
            urls = re.findall(r'https?://[^\s<>"]+|www\.[^\s<>"]+', text)
            entities["urls"] = list(set(urls))
            
            # חילוץ תאריכים
            date_patterns = [
                r'\d{1,2}/\d{1,2}/\d{2,4}',  # DD/MM/YYYY
                r'\d{1,2}-\d{1,2}-\d{2,4}',  # DD-MM-YYYY
                r'\d{1,2}\.\d{1,2}\.\d{2,4}'  # DD.MM.YYYY
            ]
            
            dates = []
            for pattern in date_patterns:
                dates.extend(re.findall(pattern, text))
            entities["dates"] = list(set(dates))
            
            # חילוץ ישויות נוספות בהתאם לשפה
            if language == "hebrew" and HEBREW_NLP_AVAILABLE:
                # שימוש בספריית hazm לניתוח עברית
                pass
            elif hasattr(nltk, 'ne_chunk'):
                # שימוש ב-NER של NLTK לזיהוי ישויות באנגלית
                try:
                    # פונקציונליות זיהוי ישויות דורשת מודלים נוספים
                    nltk.download('maxent_ne_chunker', quiet=True)
                    nltk.download('words', quiet=True)
                    
                    tokens = word_tokenize(text)
                    pos_tags = nltk.pos_tag(tokens)
                    named_entities = nltk.ne_chunk(pos_tags)
                    
                    # חילוץ ישויות לפי סוג
                    for chunk in named_entities:
                        if hasattr(chunk, 'label'):
                            entity_text = ' '.join([token for token, pos in chunk])
                            if chunk.label() == 'PERSON':
                                entities["persons"].append(entity_text)
                            elif chunk.label() == 'ORGANIZATION':
                                entities["organizations"].append(entity_text)
                            elif chunk.label() == 'GPE' or chunk.label() == 'LOCATION':
                                entities["locations"].append(entity_text)
                except Exception as e:
                    logger.debug(f"שגיאה בזיהוי ישויות: {str(e)}")
        
        except Exception as e:
            logger.error(f"שגיאה בחילוץ ישויות: {str(e)}")
        
        # הסרת כפילויות
        for entity_type in entities:
            entities[entity_type] = list(set(entities[entity_type]))
        
        return entities

    def _extract_topics(self, text: str, language: str) -> List[Dict[str, Any]]:
        """
        חילוץ נושאים מרכזיים מטקסט
        
        Args:
            text: טקסט לניתוח
            language: שפת הטקסט
            
        Returns:
            רשימת נושאים
        """
        topics = []
        
        try:
            if len(text.split()) < 50:  # לא מספיק טקסט לחילוץ נושאים
                return topics
            
            # טוקניזציה וסינון מילות עצירה
            tokens = word_tokenize(text.lower())
            
            # הגדרת מילות עצירה בהתאם לשפה
            if language == "hebrew" and HEBREW_NLP_AVAILABLE:
                stop_words = set(hazm.stopwords_list())
            else:
                stop_words = set(stopwords.words('english'))
                # הוספת מילות עצירה בעברית
                hebrew_stop_words = self._get_hebrew_stopwords()
                stop_words.update(hebrew_stop_words)
            
            filtered_tokens = [token for token in tokens if token.isalpha() and token not in stop_words]
            
            # ספירת תדירות מילים
            word_freq = collections.Counter(filtered_tokens)
            
            # הוצאת המילים הנפוצות ביותר
            most_common = word_freq.most_common(10)
            
            # בניית נושאים
            for word, freq in most_common:
                topic = {
                    "term": word,
                    "frequency": freq,
                    "score": freq / len(filtered_tokens)
                }
                topics.append(topic)
        
        except Exception as e:
            logger.error(f"שגיאה בחילוץ נושאים: {str(e)}")
        
        return topics

    def _calculate_document_similarity(self, documents: List[str]) -> np.ndarray:
        """
        חישוב מטריצת דמיון בין מסמכים
        
        Args:
            documents: רשימת תוכן טקסטואלי של מסמכים
            
        Returns:
            מטריצת דמיון בין מסמכים
        """
        if not documents or not self.vectorizer:
            return np.zeros((len(documents), len(documents)))
        
        try:
            # המרת מסמכים למטריצת TF-IDF
            tfidf_matrix = self.vectorizer.fit_transform(documents)
            
            # חישוב דמיון קוסינוס בין כל זוגות המסמכים
            similarity_matrix = cosine_similarity(tfidf_matrix)
            
            return similarity_matrix
            
        except Exception as e:
            logger.error(f"שגיאה בחישוב דמיון בין מסמכים: {str(e)}")
            return np.zeros((len(documents), len(documents)))

    def _process_generic_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קבצים כלליים (ברירת מחדל)
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {}
        
        try:
            # ניסיון לקרוא כטקסט
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text_content = f.read()
                    results["text_content"] = text_content
                    results["content_type"] = "text"
            except:
                # אם לא ניתן לקרוא כטקסט, זהו כנראה קובץ בינארי
                results["content_type"] = "binary"
                results["text_content"] = ""
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ כללי {file_path}: {str(e)}")
        
        return results

    def _process_text_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי טקסט פשוט
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "text"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
                results["line_count"] = text_content.count('\n') + 1
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ טקסט {file_path}: {str(e)}")
        
        return results

    def _process_markdown_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי Markdown
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "markdown"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
                
                # חילוץ כותרות
                headers = re.findall(r'^(#{1,6})\s+(.+?)$', text_content, re.MULTILINE)
                results["headers"] = [{"level": len(h[0]), "text": h[1]} for h in headers]
                
                # חילוץ קישורים
                links = re.findall(r'\[([^\]]+)\]\(([^)]+)\)', text_content)
                results["links"] = [{"text": link[0], "url": link[1]} for link in links]
                
                # חילוץ קטעי קוד
                code_blocks = re.findall(r'```(\w*)\n(.*?)```', text_content, re.DOTALL)
                results["code_blocks"] = [{"language": block[0], "code": block[1]} for block in code_blocks]
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ Markdown {file_path}: {str(e)}")
        
        return results

    def _process_csv_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי CSV
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "csv"}
        
        try:
            # שימוש בפנדס אם זמין
            if 'pd' in globals():
                df = pd.read_csv(file_path)
                results["row_count"] = len(df)
                results["column_count"] = len(df.columns)
                results["columns"] = list(df.columns)
                results["text_content"] = df.to_string(index=False)
                
                # הוספת סטטיסטיקה בסיסית
                try:
                    results["stats"] = {}
                    for column in df.columns:
                        if pd.api.types.is_numeric_dtype(df[column]):
                            results["stats"][column] = {
                                "min": float(df[column].min()),
                                "max": float(df[column].max()),
                                "mean": float(df[column].mean()),
                                "median": float(df[column].median())
                            }
                except Exception as stats_err:
                    logger.debug(f"שגיאה בחישוב סטטיסטיקות CSV: {str(stats_err)}")
            
            else:
                # עיבוד פשוט ללא פנדס
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    lines = f.readlines()
                    results["row_count"] = len(lines)
                    if lines:
                        results["column_count"] = len(lines[0].split(','))
                        if results["row_count"] > 0:
                            results["columns"] = [col.strip().strip('"\'') for col in lines[0].split(',')]
                    results["text_content"] = ''.join(lines)
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ CSV {file_path}: {str(e)}")
        
        return results

    def _process_json_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי JSON
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "json"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
                
                # ניסיון לפרסר את ה-JSON
                json_data = json.loads(text_content)
                
                # מידע מבני בסיסי
                if isinstance(json_data, dict):
                    results["structure"] = "object"
                    results["key_count"] = len(json_data)
                    results["top_level_keys"] = list(json_data.keys())
                elif isinstance(json_data, list):
                    results["structure"] = "array"
                    results["item_count"] = len(json_data)
                    if json_data and isinstance(json_data[0], dict):
                        # אם יש אובייקטים במערך, חלץ מפתחות נפוצים
                        all_keys = set()
                        for item in json_data:
                            if isinstance(item, dict):
                                all_keys.update(item.keys())
                        results["common_keys"] = list(all_keys)
                else:
                    results["structure"] = "primitive"
        
        except json.JSONDecodeError as json_err:
            logger.warning(f"שגיאת פירוק JSON בקובץ {file_path}: {str(json_err)}")
            results["error"] = f"Invalid JSON: {str(json_err)}"
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ JSON {file_path}: {str(e)}")
        
        return results

    def _process_yaml_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי YAML
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "yaml"}
        
        try:
            import yaml
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
                
                # ניסיון לפרסר את ה-YAML
                yaml_data = yaml.safe_load(text_content)
                
                # מידע מבני בסיסי
                if isinstance(yaml_data, dict):
                    results["structure"] = "object"
                    results["key_count"] = len(yaml_data)
                    results["top_level_keys"] = list(yaml_data.keys())
                elif isinstance(yaml_data, list):
                    results["structure"] = "array"
                    results["item_count"] = len(yaml_data)
                else:
                    results["structure"] = "primitive"
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ YAML {file_path}: {str(e)}")
            
            # אם ספריית YAML חסרה, עדיין נחזיר את התוכן
            if "yaml" not in globals():
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        results["text_content"] = f.read()
                except:
                    pass
        
        return results

    def _process_pdf_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי PDF
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "pdf"}
        
        try:
            if not DOCUMENT_CONVERSION_AVAILABLE:
                logger.warning(f"המרת PDF אינה זמינה, לא ניתן לעבד קובץ {file_path}")
                return results
            
            # שימוש ב-PyMuPDF (fitz) להמרת PDF לטקסט
            doc = fitz.open(file_path)
            
            # חילוץ מידע
            results["page_count"] = len(doc)
            results["metadata"] = doc.metadata
            
            # חילוץ תוכן טקסטואלי
            text_content = ""
            toc = []
            images = []
            
            for page_num, page in enumerate(doc):
                text_content += page.get_text()
                
                # חילוץ תמונות מהעמוד
                if self.config.get("extract_images", True):
                    for img_index, img in enumerate(page.get_images()):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_data = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # שמירת קבצי התמונות
                        image_filename = f"{os.path.basename(file_path)}_page{page_num+1}_img{img_index}.{image_ext}"
                        image_path = os.path.join(self.output_dir, image_filename)
                        
                        with open(image_path, "wb") as img_file:
                            img_file.write(image_data)
                        
                        images.append({
                            "filename": image_filename,
                            "path": image_path,
                            "page": page_num + 1,
                            "format": image_ext
                        })
            
            # חילוץ תוכן עניינים
            if doc.get_toc():
                toc = [{"level": t[0], "title": t[1], "page": t[2]} for t in doc.get_toc()]
            
            results["text_content"] = text_content
            results["toc"] = toc
            results["images"] = images
            
            doc.close()
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ PDF {file_path}: {str(e)}")
        
        return results

    def _process_docx_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי Word (DOCX)
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "docx"}
        
        try:
            if not DOCUMENT_CONVERSION_AVAILABLE:
                logger.warning(f"המרת DOCX אינה זמינה, לא ניתן לעבד קובץ {file_path}")
                return results
            
            # שימוש ב-python-docx להמרת DOCX לטקסט
            doc = docx.Document(file_path)
            
            # חילוץ טקסט
            paragraphs = [p.text for p in doc.paragraphs]
            text_content = '\n'.join(paragraphs)
            
            # חילוץ מידע
            results["paragraph_count"] = len(paragraphs)
            results["text_content"] = text_content
            
            # חילוץ כותרות
            headers = []
            for para in doc.paragraphs:
                if para.style.name.startswith('Heading'):
                    level = int(para.style.name.replace('Heading ', '')) if para.style.name != 'Heading' else 1
                    headers.append({"level": level, "text": para.text})
            
            results["headers"] = headers
            
            # חילוץ טבלאות
            tables = []
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    "id": i + 1,
                    "rows": len(table.rows),
                    "columns": len(table.rows[0].cells) if table.rows else 0,
                    "data": table_data
                })
            
            results["tables"] = tables
            
            # חילוץ תמונות
            if hasattr(doc, 'inline_shapes') and self.config.get("extract_images", True):
                images = []
                for i, shape in enumerate(doc.inline_shapes):
                    if shape.type == docx.enum.shape.WD_INLINE_SHAPE.PICTURE:
                        # שמירת התמונה
                        image_filename = f"{os.path.basename(file_path)}_img{i+1}.png"
                        image_path = os.path.join(self.output_dir, image_filename)
                        
                        try:
                            with open(image_path, "wb") as img_file:
                                img_file.write(shape.image.blob)
                            
                            images.append({
                                "filename": image_filename,
                                "path": image_path
                            })
                        except Exception as img_err:
                            logger.debug(f"שגיאה בשמירת תמונה מ-DOCX: {str(img_err)}")
                
                results["images"] = images
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ DOCX {file_path}: {str(e)}")
        
        return results

    def _process_xlsx_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי Excel (XLSX)
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "xlsx"}
        
        try:
            if not DOCUMENT_CONVERSION_AVAILABLE:
                logger.warning(f"המרת XLSX אינה זמינה, לא ניתן לעבד קובץ {file_path}")
                return results
            
            # שימוש ב-openpyxl להמרת XLSX לטקסט
            wb = openpyxl.load_workbook(file_path, data_only=True)
            
            # מידע על גיליונות
            results["sheet_count"] = len(wb.sheetnames)
            results["sheets"] = wb.sheetnames
            
            # עיבוד כל גיליון
            all_text = []
            sheets_data = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_data = {
                    "name": sheet_name,
                    "max_row": sheet.max_row,
                    "max_column": sheet.max_column,
                    "data_sample": []
                }
                
                # דגימת נתונים (עד 20 שורות ו-10 עמודות)
                for i in range(1, min(sheet.max_row + 1, 21)):
                    row_data = []
                    for j in range(1, min(sheet.max_column + 1, 11)):
                        cell_value = sheet.cell(row=i, column=j).value
                        row_data.append(str(cell_value) if cell_value is not None else "")
                        if cell_value is not None:
                            all_text.append(str(cell_value))
                    
                    sheet_data["data_sample"].append(row_data)
                
                sheets_data.append(sheet_data)
            
            results["sheets_data"] = sheets_data
            results["text_content"] = '\n'.join(all_text)
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ XLSX {file_path}: {str(e)}")
        
        return results

    def _process_pptx_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי PowerPoint (PPTX)
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "pptx"}
        
        try:
            if not DOCUMENT_CONVERSION_AVAILABLE:
                logger.warning(f"המרת PPTX אינה זמינה, לא ניתן לעבד קובץ {file_path}")
                return results
            
            # שימוש ב-python-pptx להמרת PPTX לטקסט
            presentation = pptx.Presentation(file_path)
            
            # מידע על המצגת
            results["slide_count"] = len(presentation.slides)
            
            # עיבוד כל שקופית
            all_text = []
            slides_data = []
            
            for i, slide in enumerate(presentation.slides):
                slide_data = {
                    "index": i + 1,
                    "shapes": len(slide.shapes),
                    "text": []
                }
                
                # חילוץ טקסט מהשקופית
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                        all_text.append(shape.text)
                
                slide_data["text"] = slide_text
                slides_data.append(slide_data)
            
            results["slides"] = slides_data
            results["text_content"] = '\n'.join(all_text)
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ PPTX {file_path}: {str(e)}")
        
        return results

    def _process_python_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי Python
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "python"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
            
            # חילוץ יבוא מודולים
            import_pattern = r'(?:import|from)\s+([\w.]+)'
            imports = re.findall(import_pattern, text_content)
            results["imports"] = list(set(imports))
            
            # חילוץ פונקציות
            function_pattern = r'def\s+(\w+)\s*\(([^)]*)\)'
            functions = re.findall(function_pattern, text_content)
            results["functions"] = [{"name": f[0], "params": f[1].strip()} for f in functions]
            
            # חילוץ מחלקות
            class_pattern = r'class\s+(\w+)(?:\(([^)]*)\))?'
            classes = re.findall(class_pattern, text_content)
            results["classes"] = [{"name": c[0], "inherits": c[1].strip()} for c in classes]
            
            # חילוץ מחרוזות docstring
            docstring_pattern = r'"{3}(.*?)"{3}|\'"{3}(.*?)\'"{3}'
            docstrings = re.findall(docstring_pattern, text_content, re.DOTALL)
            results["docstrings"] = [d[0] or d[1] for d in docstrings]
            
            # חילוץ הערות כלליות
            comment_pattern = r'#\s*(.*)'
            comments = re.findall(comment_pattern, text_content)
            results["comments"] = comments
            
            # חילוץ מידע על מבנה הקוד
            try:
                ast_tree = ast.parse(text_content)
                
                # ספירת סוגי הצהרות
                stmt_types = collections.Counter()
                for node in ast.walk(ast_tree):
                    stmt_types[type(node).__name__] += 1
                
                results["syntax_elements"] = {k: v for k, v in stmt_types.most_common()}
            except SyntaxError as syntax_err:
                results["syntax_error"] = str(syntax_err)
                results["syntax_error_line"] = syntax_err.lineno
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ Python {file_path}: {str(e)}")
        
        return results

    def _process_javascript_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי JavaScript
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "javascript"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
            
            # חילוץ יבוא מודולים
            import_pattern = r'(?:import|require)\s*\(?\s*[\'"]([\w./-]+)[\'"]'
            imports = re.findall(import_pattern, text_content)
            results["imports"] = list(set(imports))
            
            # חילוץ פונקציות
            function_pattern = r'function\s+(\w+)\s*\(([^)]*)\)'
            functions = re.findall(function_pattern, text_content)
            
            # כולל פונקציות חץ
            arrow_function_pattern = r'(?:const|let|var)\s+(\w+)\s*=\s*\(([^)]*)\)\s*=>'
            arrow_functions = re.findall(arrow_function_pattern, text_content)
            
            all_functions = [{"name": f[0], "params": f[1].strip(), "type": "function"} for f in functions]
            all_functions.extend([{"name": f[0], "params": f[1].strip(), "type": "arrow"} for f in arrow_functions])
            
            results["functions"] = all_functions
            
            # חילוץ מחלקות
            class_pattern = r'class\s+(\w+)(?:\s+extends\s+(\w+))?'
            classes = re.findall(class_pattern, text_content)
            results["classes"] = [{"name": c[0], "extends": c[1] if c[1] else ""} for c in classes]
            
            # חילוץ הערות JSDoc
            jsdoc_pattern = r'/\*\*(.*?)\*/'
            jsdocs = re.findall(jsdoc_pattern, text_content, re.DOTALL)
            results["jsdocs"] = [doc.strip() for doc in jsdocs]
            
            # חילוץ הערות כלליות
            comment_pattern = r'//\s*(.*)'
            comments = re.findall(comment_pattern, text_content)
            results["comments"] = comments
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ JavaScript {file_path}: {str(e)}")
        
        return results

    def _process_html_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי HTML
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "html"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
            
            # חילוץ תגיות Head
            head_pattern = r'<head.*?>(.*?)</head>'
            head_match = re.search(head_pattern, text_content, re.DOTALL)
            if head_match:
                head_content = head_match.group(1)
                
                # חילוץ כותרת
                title_match = re.search(r'<title>(.*?)</title>', head_content)
                if title_match:
                    results["title"] = title_match.group(1)
                
                # חילוץ meta
                meta_tags = re.findall(r'<meta\s+([^>]+)>', head_content)
                meta_data = []
                
                for meta in meta_tags:
                    meta_attrs = {}
                    name_match = re.search(r'name\s*=\s*["\']([^"\']+)["\']', meta)
                    if name_match:
                        meta_attrs["name"] = name_match.group(1)
                    
                    content_match = re.search(r'content\s*=\s*["\']([^"\']+)["\']', meta)
                    if content_match:
                        meta_attrs["content"] = content_match.group(1)
                    
                    if meta_attrs:
                        meta_data.append(meta_attrs)
                
                results["meta_tags"] = meta_data
            
            # חילוץ תגיות Body
            body_pattern = r'<body.*?>(.*?)</body>'
            body_match = re.search(body_pattern, text_content, re.DOTALL)
            if body_match:
                body_content = body_match.group(1)
                
                # חילוץ טקסט (הסרת HTML)
                text_only = re.sub(r'<[^>]+>', ' ', body_content)
                text_only = re.sub(r'\s+', ' ', text_only).strip()
                results["body_text"] = text_only
                
                # ספירת תגיות
                tag_counts = collections.Counter(re.findall(r'<(\w+)[^>]*>', body_content))
                results["tag_counts"] = dict(tag_counts.most_common(10))
                
                # חילוץ קישורים
                links = re.findall(r'<a\s+[^>]*href\s*=\s*["\']([^"\']+)["\'][^>]*>(.*?)</a>', body_content, re.DOTALL)
                results["links"] = [{"url": link[0], "text": re.sub(r'<[^>]+>', '', link[1]).strip()} for link in links]
                
                # חילוץ תמונות
                images = re.findall(r'<img\s+[^>]*src\s*=\s*["\']([^"\']+)["\'][^>]*>', body_content)
                results["images"] = images
            
            # ניתוח סקריפטים
            scripts = re.findall(r'<script[^>]*>(.*?)</script>', text_content, re.DOTALL)
            results["scripts"] = [script.strip() for script in scripts]
            
            # ניתוח סגנונות
            styles = re.findall(r'<style[^>]*>(.*?)</style>', text_content, re.DOTALL)
            results["styles"] = [style.strip() for style in styles]
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ HTML {file_path}: {str(e)}")
        
        return results

    def _process_css_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי CSS
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "css"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
            
            # חילוץ מסננים
            selector_pattern = r'([^{]+){([^}]*)}'
            selectors = re.findall(selector_pattern, text_content)
            
            # עיבוד מסננים
            processed_selectors = []
            for selector, properties in selectors:
                selector = selector.strip()
                props = {}
                
                # עיבוד תכונות
                prop_pattern = r'([\w-]+)\s*:\s*([^;]+);?'
                for prop, value in re.findall(prop_pattern, properties):
                    props[prop.strip()] = value.strip()
                
                processed_selectors.append({
                    "selector": selector,
                    "properties": props,
                    "property_count": len(props)
                })
            
            results["selectors"] = processed_selectors
            results["selector_count"] = len(processed_selectors)
            
            # ניתוח צבעים
            color_pattern = r'#[0-9a-fA-F]{3,6}|rgba?\([^)]+\)|hsla?\([^)]+\)'
            colors = re.findall(color_pattern, text_content)
            results["colors"] = list(set(colors))
            
            # ניתוח media queries
            media_queries = re.findall(r'@media\s+([^{]+){([^}]*)(?:}[^{]*{[^}]*})*', text_content)
            results["media_queries"] = [query[0].strip() for query in media_queries]
            
            # ניתוח keyframes
            keyframes = re.findall(r'@keyframes\s+([^{]+){([^}]*)(?:}[^{]*{[^}]*})*', text_content)
            results["keyframes"] = [keyframe[0].strip() for keyframe in keyframes]
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ CSS {file_path}: {str(e)}")
        
        return results

    def _process_xml_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי XML
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "xml"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
            
            # בדיקת הסכמה
            xml_schema = None
            schema_match = re.search(r'<\?xml-stylesheet[^>]*href\s*=\s*["\']([^"\']+)["\']', text_content)
            if schema_match:
                xml_schema = schema_match.group(1)
                results["schema"] = xml_schema
            
            # חילוץ רשימת תגיות
            tag_pattern = r'<([a-zA-Z0-9_:-]+)[^>]*>'
            tags = re.findall(tag_pattern, text_content)
            tag_counts = collections.Counter(tags)
            results["tag_counts"] = dict(tag_counts.most_common(10))
            
            # ניסיון להבין את המבנה ההיררכי
            root_tag_match = re.search(r'<([a-zA-Z0-9_:-]+)[^>]*>', text_content)
            if root_tag_match:
                results["root_tag"] = root_tag_match.group(1)
            
            # חילוץ ערכי תגיות
            value_pattern = r'<([a-zA-Z0-9_:-]+)[^>]*>(.*?)</\1>'
            values = re.findall(value_pattern, text_content, re.DOTALL)
            
            # שמירת ערכים מאובחנים (עד 10)
            extracted_values = {}
            for i, (tag, value) in enumerate(values[:10]):
                # ניקוי הערך מתגיות נוספות
                clean_value = re.sub(r'<[^>]+>', '', value).strip()
                if clean_value:
                    extracted_values[tag] = clean_value
            
            results["extracted_values"] = extracted_values
            
            # חילוץ תכונות
            attr_pattern = r'<[^>]+\s+([a-zA-Z0-9_:-]+)\s*=\s*["\']([^"\']*)["\']'
            attrs = re.findall(attr_pattern, text_content)
            attr_counts = collections.Counter([attr[0] for attr in attrs])
            results["attribute_counts"] = dict(attr_counts.most_common(10))
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ XML {file_path}: {str(e)}")
        
        return results

    def _process_image_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי תמונה
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "image"}
        
        try:
            if not IMAGE_PROCESSING_AVAILABLE:
                logger.warning(f"עיבוד תמונה אינו זמין, לא ניתן לנתח קובץ {file_path}")
                return results
            
            # קריאת התמונה עם PIL
            image = Image.open(file_path)
            
            # מידע בסיסי
            results["format"] = image.format
            results["mode"] = image.mode
            results["size"] = {"width": image.width, "height": image.height}
            
            # חילוץ מידע EXIF אם קיים
            exif_data = {}
            if hasattr(image, '_getexif') and image._getexif():
                exif = image._getexif()
                for tag_id, value in exif.items():
                    tag_name = f"EXIF_{tag_id}"
                    # המרת ערכים לפורמט JSON-serializable
                    if isinstance(value, bytes):
                        value = str(value)
                    exif_data[tag_name] = value
            
            results["exif"] = exif_data
            
            # זיהוי טקסט בתמונה (OCR) אם זמין
            if self.config.get("use_ocr", True) and hasattr(pytesseract, 'image_to_string'):
                try:
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():  # אם נמצא טקסט
                        results["ocr_text"] = ocr_text.strip()
                        results["text_content"] = ocr_text.strip()  # לתאימות עם ניתוח טקסט
                except Exception as ocr_err:
                    logger.debug(f"שגיאה בזיהוי טקסט בתמונה: {str(ocr_err)}")
            
            # חילוץ צבעים דומיננטיים
            if self.config.get("extract_colors", True):
                try:
                    # המרה למצב RGB אם צריך
                    if image.mode != 'RGB':
                        image = image.convert('RGB')
                    
                    # דגימת פיקסלים
                    pixels = image.resize((50, 50)).getdata()
                    colors = collections.Counter(pixels)
                    most_common = colors.most_common(5)
                    
                    dominant_colors = []
                    for color, count in most_common:
                        hex_color = "#{:02x}{:02x}{:02x}".format(*color)
                        dominant_colors.append({
                            "rgb": color,
                            "hex": hex_color,
                            "count": count
                        })
                    
                    results["dominant_colors"] = dominant_colors
                except Exception as color_err:
                    logger.debug(f"שגיאה בחילוץ צבעים דומיננטיים: {str(color_err)}")
            
            # המרה לבסיס 64 לתצוגה
            if self.config.get("include_preview", True):
                img_preview = image.copy()
                max_size = (400, 400)
                img_preview.thumbnail(max_size)
                
                buffer = BytesIO()
                img_preview.save(buffer, format=image.format or "PNG")
                img_data = base64.b64encode(buffer.getvalue()).decode('utf-8')
                
                results["preview"] = f"data:image/{image.format.lower() if image.format else 'png'};base64,{img_data}"
            
            # סגירת התמונה
            image.close()
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ תמונה {file_path}: {str(e)}")
        
        return results

    def _process_svg_file(self, file_path: str) -> Dict[str, Any]:
        """
        מעבד קובצי SVG
        
        Args:
            file_path: נתיב הקובץ
            
        Returns:
            מילון עם תוצאות העיבוד
        """
        results = {"content_type": "svg"}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
                results["text_content"] = text_content
            
            # מידע בסיסי
            svg_tag = re.search(r'<svg[^>]*>', text_content)
            if svg_tag:
                # ניסיון לחלץ גודל
                width_match = re.search(r'width\s*=\s*["\']([^"\']+)["\']', svg_tag.group(0))
                height_match = re.search(r'height\s*=\s*["\']([^"\']+)["\']', svg_tag.group(0))
                
                if width_match and height_match:
                    results["size"] = {
                        "width": width_match.group(1),
                        "height": height_match.group(1)
                    }
                
                # חילוץ viewBox
                viewbox_match = re.search(r'viewBox\s*=\s*["\']([^"\']+)["\']', svg_tag.group(0))
                if viewbox_match:
                    results["viewBox"] = viewbox_match.group(1)
            
            # חילוץ אלמנטים
            elements = {}
            for tag in ['path', 'rect', 'circle', 'ellipse', 'line', 'polyline', 'polygon', 'text', 'g', 'defs']:
                count = len(re.findall(f'<{tag}[^>]*>', text_content))
                if count > 0:
                    elements[tag] = count
            
            results["elements"] = elements
            
            # חילוץ טקסט בתוך ה-SVG
            text_blocks = re.findall(r'<text[^>]*>(.*?)</text>', text_content, re.DOTALL)
            if text_blocks:
                # ניקוי הטקסט מתגיות פנימיות
                cleaned_text = []
                for block in text_blocks:
                    clean_block = re.sub(r'<[^>]+>', '', block).strip()
                    if clean_block:
                        cleaned_text.append(clean_block)
                
                results["svg_text"] = cleaned_text
            
            # חילוץ צבעים
            color_pattern = r'(?:fill|stroke)\s*=\s*["\']([^"\']+)["\']'
            colors = re.findall(color_pattern, text_content)
            results["colors"] = list(set(colors))
        
        except Exception as e:
            logger.error(f"שגיאה בעיבוד קובץ SVG {file_path}: {str(e)}")
        
        return results

    def extract_text_from_document(self, file_path: str) -> str:
        """
        חילוץ טקסט ממסמך בכל פורמט
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            טקסט מחולץ מהמסמך
        """
        # ניתוח המסמך
        results = self.analyze_document(file_path)
        
        # החזרת הטקסט אם קיים
        if "text_content" in results:
            return results["text_content"]
        
        return ""

    def find_similar_documents(self, base_file: str, files_to_compare: List[str]) -> List[Dict[str, Any]]:
        """
        חיפוש מסמכים דומים למסמך בסיס
        
        Args:
            base_file: נתיב למסמך בסיס
            files_to_compare: רשימת נתיבים למסמכים להשוואה
            
        Returns:
            רשימת מסמכים דומים מדורגת לפי דמיון
        """
        if not NLP_AVAILABLE:
            logger.warning("ניתוח NLP אינו זמין, לא ניתן לחשב דמיון בין מסמכים")
            return []
        
        try:
            # חילוץ טקסט מכל המסמכים
            base_text = self.extract_text_from_document(base_file)
            
            if not base_text:
                logger.warning(f"לא ניתן לחלץ טקסט ממסמך הבסיס {base_file}")
                return []
            
            texts_to_compare = []
            for file_path in files_to_compare:
                text = self.extract_text_from_document(file_path)
                texts_to_compare.append((file_path, text))
            
            # סינון מסמכים ללא טקסט
            valid_texts = [(path, text) for path, text in texts_to_compare if text]
            
            if not valid_texts:
                logger.warning("לא נמצאו מסמכים עם טקסט להשוואה")
                return []
            
            # חישוב דמיון
            similarities = []
            
            for file_path, text in valid_texts:
                # חישוב דמיון קוסינוס
                vectorizer = TfidfVectorizer()
                tfidf_matrix = vectorizer.fit_transform([base_text, text])
                similarity = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])[0][0]
                
                similarities.append({
                    "file_path": file_path,
                    "similarity": float(similarity)
                })
            
            # מיון לפי דמיון
            similarities.sort(key=lambda x: x["similarity"], reverse=True)
            
            # סינון תוצאות מתחת לסף
            return [s for s in similarities if s["similarity"] >= self.min_similarity]
            
        except Exception as e:
            logger.error(f"שגיאה בחיפוש מסמכים דומים: {str(e)}")
            return []

    def create_document_summary(self, file_path: str, max_length: int = 500) -> Dict[str, Any]:
        """
        יצירת סיכום מסמך
        
        Args:
            file_path: נתיב למסמך
            max_length: אורך מקסימלי לסיכום
            
        Returns:
            מילון עם סיכום המסמך
        """
        if not NLP_AVAILABLE:
            logger.warning("ניתוח NLP אינו זמין, לא ניתן ליצור סיכום מסמך")
            return {"status": "error", "error": "ניתוח NLP אינו זמין"}
        
        try:
            # ניתוח המסמך
            results = self.analyze_document(file_path)
            
            if results.get("status") != "success":
                return results
            
            # קבלת הטקסט
            text_content = results.get("text_content", "")
            
            if not text_content:
                return {"status": "error", "error": "לא נמצא תוכן טקסט במסמך"}
            
            # סיכום הטקסט
            summary = ""
            if len(text_content.split()) > 100:  # רק אם יש מספיק טקסט
                try:
                    # גודל הסיכום תלוי באורך הטקסט המקורי
                    ratio = min(max_length / len(text_content), 0.5)
                    summary = summarize(text_content, ratio=ratio)
                except:
                    # אם הסיכום האוטומטי נכשל, יצירת סיכום פשוט
                    sentences = sent_tokenize(text_content)
                    summary = ' '.join(sentences[:5])  # 5 משפטים ראשונים
            
            # מילות מפתח
            keywords_list = []
            try:
                keywords_text = keywords(text_content, ratio=self.keywords_ratio)
                keywords_list = keywords_text.split('\n') if keywords_text else []
            except:
                # אם חילוץ מילות מפתח נכשל, נחלץ מילים נפוצות
                words = word_tokenize(text_content.lower())
                if self.stop_words:
                    words = [w for w in words if w.isalpha() and w not in self.stop_words]
                word_freq = collections.Counter(words)
                keywords_list = [word for word, _ in word_freq.most_common(10)]
            
            # יצירת סיכום מסמך
            document_summary = {
                "status": "success",
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "mime_type": results.get("mime_type", ""),
                "summary": summary,
                "keywords": keywords_list,
                "word_count": results.get("word_count", 0),
                "char_count": results.get("char_count", 0),
                "language": results.get("language", "")
            }
            
            # הוספת מידע ספציפי לסוג המסמך
            if results.get("content_type") == "pdf":
                document_summary["page_count"] = results.get("page_count", 0)
                document_summary["toc"] = results.get("toc", [])
            elif results.get("content_type") == "docx":
                document_summary["paragraph_count"] = results.get("paragraph_count", 0)
                document_summary["headers"] = results.get("headers", [])
            
            return document_summary
            
        except Exception as e:
            logger.error(f"שגיאה ביצירת סיכום מסמך {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def scan_document_for_keywords(self, file_path: str, keywords: List[str]) -> Dict[str, Any]:
        """
        סריקת מסמך למציאת מילות מפתח
        
        Args:
            file_path: נתיב למסמך
            keywords: רשימת מילות מפתח לחיפוש
            
        Returns:
            מילון עם תוצאות החיפוש
        """
        try:
            # חילוץ טקסט מהמסמך
            text_content = self.extract_text_from_document(file_path)
            
            if not text_content:
                return {
                    "status": "error", 
                    "error": "לא ניתן לחלץ טקסט מהמסמך",
                    "file_path": file_path
                }
            
            # חיפוש מילות מפתח
            results = {"status": "success", "file_path": file_path, "matches": {}}
            
            for keyword in keywords:
                # יצירת דפוס חיפוש
                pattern = re.compile(r'\b' + re.escape(keyword) + r'\b', re.IGNORECASE)
                matches = pattern.finditer(text_content)
                
                # איסוף מידע על התאמות
                occurrences = []
                for match in matches:
                    # חילוץ קונטקסט
                    start = max(0, match.start() - 50)
                    end = min(len(text_content), match.end() + 50)
                    context = text_content[start:end]
                    
                    # הוספת מיקום
                    line_number = text_content[:match.start()].count('\n') + 1
                    pos = match.start() - text_content[:match.start()].rfind('\n') if text_content[:match.start()].rfind('\n') >= 0 else match.start()
                    
                    occurrences.append({
                        "position": match.start(),
                        "line": line_number,
                        "column": pos,
                        "context": context
                    })
                
                if occurrences:
                    results["matches"][keyword] = {
                        "count": len(occurrences),
                        "occurrences": occurrences
                    }
            
            # מידע סיכומי
            total_matches = sum(data["count"] for data in results["matches"].values())
            results["total_matches"] = total_matches
            results["matched_keywords"] = list(results["matches"].keys())
            
            return results
            
        except Exception as e:
            logger.error(f"שגיאה בסריקת מסמך {file_path} למציאת מילות מפתח: {str(e)}")
            return {"status": "error", "error": str(e), "file_path": file_path}

    def convert_document(self, file_path: str, output_format: str) -> Dict[str, Any]:
        """
        המרת מסמך לפורמט אחר
        
        Args:
            file_path: נתיב למסמך
            output_format: פורמט הפלט
            
        Returns:
            מילון עם תוצאות ההמרה
        """
        if not DOCUMENT_CONVERSION_AVAILABLE:
            logger.warning("המרת מסמכים אינה זמינה")
            return {"status": "error", "error": "המרת מסמכים אינה זמינה"}
        
        try:
            # זיהוי סוג הקובץ
            mime_type = self._detect_mime_type(file_path)
            
            # יצירת שם קובץ פלט
            output_filename = os.path.splitext(os.path.basename(file_path))[0] + "." + output_format.lower()
            output_path = os.path.join(self.output_dir, output_filename)
            
            # המרה בהתאם לפורמט המקור והיעד
            if mime_type == 'application/pdf' and output_format.lower() == 'txt':
                # המרת PDF לטקסט
                doc = fitz.open(file_path)
                with open(output_path, 'w', encoding='utf-8') as f:
                    for page in doc:
                        f.write(page.get_text())
                doc.close()
            
            elif mime_type == 'application/pdf' and output_format.lower() == 'html':
                # המרת PDF ל-HTML
                doc = fitz.open(file_path)
                html_content = ["<html><head><title>Converted PDF</title></head><body>"]
                
                for page in doc:
                    html_content.append(f"<div class='page' id='page-{page.number+1}'>")
                    html_content.append(f"<h2>Page {page.number+1}</h2>")
                    html_content.append(f"<div class='content'>{page.get_text('html')}</div>")
                    html_content.append("</div>")
                
                html_content.append("</body></html>")
                
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write("\n".join(html_content))
                
                doc.close()
            
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' and output_format.lower() == 'txt':
                # המרת DOCX לטקסט
                doc = docx.Document(file_path)
                with open(output_path, 'w', encoding='utf-8') as f:
                    for para in doc.paragraphs:
                        f.write(para.text + '\n')
            
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' and output_format.lower() == 'html':
                # המרת DOCX ל-HTML
                doc = docx.Document(file_path)
                html_content = ["<html><head><title>Converted DOCX</title></head><body>"]
                
                for para in doc.paragraphs:
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name.replace('Heading ', '')) if para.style.name != 'Heading' else 1
                        html_content.append(f"<h{level}>{para.text}</h{level}>")
                    else:
                        html_content.append(f"<p>{para.text}</p>")
                
                html_content.append("</body></html>")
                
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write("\n".join(html_content))
            
            else:
                return {"status": "error", "error": f"המרה מ-{mime_type} ל-{output_format} אינה נתמכת"}
            
            return {
                "status": "success",
                "source_file": file_path,
                "output_file": output_path,
                "output_format": output_format.lower()
            }
            
        except Exception as e:
            logger.error(f"שגיאה בהמרת מסמך {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def extract_structured_data(self, file_path: str) -> Dict[str, Any]:
        """
        חילוץ מידע מובנה ממסמך
        
        Args:
            file_path: נתיב למסמך
            
        Returns:
            מילון עם המידע המובנה
        """
        try:
            # זיהוי סוג הקובץ
            mime_type = self._detect_mime_type(file_path)
            
            # חילוץ מידע בהתאם לסוג הקובץ
            if mime_type == 'application/pdf':
                # חילוץ מידע מובנה מ-PDF
                return self._extract_structured_from_pdf(file_path)
            
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                # חילוץ מידע מובנה מ-DOCX
                return self._extract_structured_from_docx(file_path)
            
            elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                # חילוץ מידע מובנה מ-XLSX
                return self._extract_structured_from_xlsx(file_path)
            
            elif mime_type == 'text/csv':
                # חילוץ מידע מובנה מ-CSV
                return self._extract_structured_from_csv(file_path)
            
            elif mime_type == 'application/json':
                # חילוץ מידע מובנה מ-JSON
                return self._extract_structured_from_json(file_path)
            
            else:
                # ניסיון חילוץ כללי
                return self._extract_structured_generic(file_path)
            
        except Exception as e:
            logger.error(f"שגיאה בחילוץ מידע מובנה ממסמך {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def _extract_structured_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        חילוץ מידע מובנה מקובץ PDF
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            מילון עם המידע המובנה
        """
        if not DOCUMENT_CONVERSION_AVAILABLE:
            return {"status": "error", "error": "המרת PDF אינה זמינה"}
        
        structured_data = {"status": "success", "file_path": file_path, "file_type": "PDF"}
        
        try:
            doc = fitz.open(file_path)
            
            # חילוץ מטא-דאטה
            structured_data["metadata"] = doc.metadata
            
            # חילוץ תוכן עניינים
            if doc.get_toc():
                structured_data["toc"] = [{"level": t[0], "title": t[1], "page": t[2]} for t in doc.get_toc()]
            
            # חילוץ טבלאות (פשוט - זיהוי טקסט בתבנית של טבלה)
            tables = []
            for page_num, page in enumerate(doc):
                # חילוץ טקסט בפורמט בלוקים
                blocks = page.get_text("blocks")
                
                # זיהוי בלוקים שעשויים להיות טבלאות
                for block in blocks:
                    text = block[4]
                    lines = text.strip().split('\n')
                    
                    # בדיקה אם הבלוק מכיל מבנה של טבלה
                    if len(lines) > 2:
                        # בדיקה אם יש אותו מספר עמודות בכל שורה
                        columns = [len(line.split()) for line in lines]
                        if len(set(columns)) <= 2 and min(columns) >= 2:
                            # כנראה טבלה
                            table_data = [line.split() for line in lines]
                            tables.append({
                                "page": page_num + 1,
                                "rows": len(table_data),
                                "columns": max(columns),
                                "data": table_data
                            })
            
            structured_data["tables"] = tables
            
            # חילוץ מידע מובנה (לדוגמה: שדות טופס)
            form_fields = []
            if doc.has_links():
                for page_num, page in enumerate(doc):
                    for field in page.widgets():
                        if field.field_type_string:
                            form_fields.append({
                                "page": page_num + 1,
                                "type": field.field_type_string,
                                "name": field.field_name,
                                "value": field.field_value,
                                "rect": list(field.rect)
                            })
            
            structured_data["form_fields"] = form_fields
            
            doc.close()
            
            return structured_data
            
        except Exception as e:
            logger.error(f"שגיאה בחילוץ מידע מובנה מ-PDF {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def _extract_structured_from_docx(self, file_path: str) -> Dict[str, Any]:
        """
        חילוץ מידע מובנה מקובץ DOCX
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            מילון עם המידע המובנה
        """
        if not DOCUMENT_CONVERSION_AVAILABLE:
            return {"status": "error", "error": "המרת DOCX אינה זמינה"}
        
        structured_data = {"status": "success", "file_path": file_path, "file_type": "DOCX"}
        
        try:
            doc = docx.Document(file_path)
            
            # חילוץ מאפייני המסמך
            core_properties = {}
            if hasattr(doc, 'core_properties'):
                props = doc.core_properties
                for prop_name in dir(props):
                    if not prop_name.startswith('_') and not callable(getattr(props, prop_name)):
                        value = getattr(props, prop_name)
                        if value is not None:
                            core_properties[prop_name] = str(value)
            
            structured_data["core_properties"] = core_properties
            
            # חילוץ פסקאות
            paragraphs = []
            for p in doc.paragraphs:
                if p.text.strip():
                    paragraphs.append({
                        "text": p.text,
                        "style": p.style.name
                    })
            
            structured_data["paragraphs"] = paragraphs
            
            # חילוץ כותרות
            headings = []
            for para in doc.paragraphs:
                if para.style.name.startswith('Heading'):
                    level = int(para.style.name.replace('Heading ', '')) if para.style.name != 'Heading' else 1
                    headings.append({
                        "level": level,
                        "text": para.text
                    })
            
            structured_data["headings"] = headings
            
            # חילוץ טבלאות
            tables = []
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                tables.append({
                    "id": i + 1,
                    "rows": len(table.rows),
                    "columns": len(table.rows[0].cells) if table.rows else 0,
                    "data": table_data
                })
            
            structured_data["tables"] = tables
            
            return structured_data
            
        except Exception as e:
            logger.error(f"שגיאה בחילוץ מידע מובנה מ-DOCX {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def _extract_structured_from_xlsx(self, file_path: str) -> Dict[str, Any]:
        """
        חילוץ מידע מובנה מקובץ XLSX
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            מילון עם המידע המובנה
        """
        if not DOCUMENT_CONVERSION_AVAILABLE:
            return {"status": "error", "error": "המרת XLSX אינה זמינה"}
        
        structured_data = {"status": "success", "file_path": file_path, "file_type": "XLSX"}
        
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            
            # מידע כללי
            structured_data["sheet_names"] = wb.sheetnames
            
            # חילוץ נתונים מכל גיליון
            sheets_data = {}
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # נתונים בסיסיים
                sheet_info = {
                    "max_row": sheet.max_row,
                    "max_column": sheet.max_column,
                    "data": [],
                    "headers": []
                }
                
                # זיהוי כותרות (שורה ראשונה)
                if sheet.max_row > 0 and sheet.max_column > 0:
                    headers = []
                    for col in range(1, sheet.max_column + 1):
                        cell_value = sheet.cell(row=1, column=col).value
                        headers.append(str(cell_value) if cell_value is not None else "")
                    
                    sheet_info["headers"] = headers
                
                # חילוץ נתונים (עד 1000 שורות לביצועים טובים)
                max_rows_to_read = min(sheet.max_row, 1000)
                for row in range(2, max_rows_to_read + 1):  # מתחיל משורה 2 (אחרי הכותרות)
                    row_data = {}
                    for col in range(1, sheet.max_column + 1):
                        column_name = sheet_info["headers"][col-1] if col-1 < len(sheet_info["headers"]) else f"Column {col}"
                        cell_value = sheet.cell(row=row, column=col).value
                        row_data[column_name] = cell_value
                    
                    sheet_info["data"].append(row_data)
                
                # הוספת נתוני הגיליון
                sheets_data[sheet_name] = sheet_info
            
            structured_data["sheets"] = sheets_data
            
            # חילוץ מאפייני המסמך
            if hasattr(wb, 'properties'):
                properties = {}
                for prop_name in dir(wb.properties):
                    if not prop_name.startswith('_') and not callable(getattr(wb.properties, prop_name)):
                        value = getattr(wb.properties, prop_name)
                        if value is not None:
                            properties[prop_name] = str(value)
                
                structured_data["properties"] = properties
            
            return structured_data
            
        except Exception as e:
            logger.error(f"שגיאה בחילוץ מידע מובנה מ-XLSX {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def _extract_structured_from_csv(self, file_path: str) -> Dict[str, Any]:
        """
        חילוץ מידע מובנה מקובץ CSV
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            מילון עם המידע המובנה
        """
        structured_data = {"status": "success", "file_path": file_path, "file_type": "CSV"}
        
        try:
            # שימוש בפנדס אם זמין
            if 'pd' in globals():
                df = pd.read_csv(file_path)
                
                # נתונים בסיסיים
                structured_data["row_count"] = len(df)
                structured_data["column_count"] = len(df.columns)
                structured_data["columns"] = list(df.columns)
                
                # דגימת נתונים (עד 100 שורות)
                structured_data["data"] = df.head(100).to_dict(orient='records')
                
                # סטטיסטיקה בסיסית
                stats = {}
                for column in df.columns:
                    if pd.api.types.is_numeric_dtype(df[column]):
                        column_stats = {
                            "min": float(df[column].min()),
                            "max": float(df[column].max()),
                            "mean": float(df[column].mean()),
                            "median": float(df[column].median()),
                            "std": float(df[column].std())
                        }
                        stats[column] = column_stats
                
                structured_data["statistics"] = stats
                
                # זיהוי סוגי נתונים
                column_types = {}
                for column in df.columns:
                    dtype = str(df[column].dtype)
                    column_types[column] = dtype
                
                structured_data["column_types"] = column_types
            
            else:
                # חילוץ פשוט אם פנדס לא זמין
                import csv
                
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    csv_reader = csv.reader(f)
                    rows = list(csv_reader)
                    
                    if rows:
                        headers = rows[0]
                        data = []
                        
                        for row in rows[1:101]:  # עד 100 שורות
                            row_data = {}
                            for i, value in enumerate(row):
                                if i < len(headers):
                                    row_data[headers[i]] = value
                            data.append(row_data)
                        
                        structured_data["row_count"] = len(rows) - 1
                        structured_data["column_count"] = len(headers)
                        structured_data["columns"] = headers
                        structured_data["data"] = data
            
            return structured_data
            
        except Exception as e:
            logger.error(f"שגיאה בחילוץ מידע מובנה מ-CSV {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def _extract_structured_from_json(self, file_path: str) -> Dict[str, Any]:
        """
        חילוץ מידע מובנה מקובץ JSON
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            מילון עם המידע המובנה
        """
        structured_data = {"status": "success", "file_path": file_path, "file_type": "JSON"}
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                json_data = json.load(f)
            
            # הוספת הנתונים
            structured_data["data"] = json_data
            
            # ניתוח המבנה
            if isinstance(json_data, dict):
                structured_data["structure_type"] = "object"
                structured_data["top_level_keys"] = list(json_data.keys())
                structured_data["key_count"] = len(json_data)
                
                # ניתוח שדות מרכזיים
                key_types = {}
                for key, value in json_data.items():
                    if isinstance(value, (str, int, float, bool)):
                        key_types[key] = type(value).__name__
                    elif isinstance(value, list):
                        key_types[key] = f"array[{len(value)}]"
                    elif isinstance(value, dict):
                        key_types[key] = f"object[{len(value)}]"
                    else:
                        key_types[key] = "unknown"
                
                structured_data["key_types"] = key_types
                
            elif isinstance(json_data, list):
                structured_data["structure_type"] = "array"
                structured_data["item_count"] = len(json_data)
                
                # בדיקה אם יש מבנה אחיד
                if json_data and all(isinstance(item, dict) for item in json_data):
                    # איסוף כל המפתחות
                    all_keys = set()
                    for item in json_data:
                        all_keys.update(item.keys())
                    
                    structured_data["common_keys"] = list(all_keys)
                    
                    # חישוב כמה פעמים מופיע כל מפתח
                    key_frequencies = {}
                    for key in all_keys:
                        frequency = sum(1 for item in json_data if key in item)
                        key_frequencies[key] = frequency
                    
                    structured_data["key_frequencies"] = key_frequencies
            else:
                structured_data["structure_type"] = "primitive"
            
            return structured_data
            
        except Exception as e:
            logger.error(f"שגיאה בחילוץ מידע מובנה מ-JSON {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def _extract_structured_generic(self, file_path: str) -> Dict[str, Any]:
        """
        חילוץ מידע מובנה כללי ממסמך
        
        Args:
            file_path: נתיב לקובץ
            
        Returns:
            מילון עם המידע המובנה
        """
        mime_type = self._detect_mime_type(file_path)
        structured_data = {"status": "success", "file_path": file_path, "file_type": mime_type}
        
        try:
            # קריאת הקובץ כטקסט
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text_content = f.read()
            except:
                return {"status": "error", "error": "לא ניתן לקרוא את הקובץ כטקסט"}
            
            # חילוץ נתונים בסיסיים
            structured_data["file_size"] = os.path.getsize(file_path)
            structured_data["line_count"] = text_content.count('\n') + 1
            structured_data["word_count"] = len(text_content.split())
            structured_data["char_count"] = len(text_content)
            
            # חילוץ כתובות אימייל
            emails = re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', text_content)
            if emails:
                structured_data["emails"] = list(set(emails))
            
            # חילוץ כתובות URL
            urls = re.findall(r'https?://[^\s<>"]+|www\.[^\s<>"]+', text_content)
            if urls:
                structured_data["urls"] = list(set(urls))
            
            # חילוץ מספרי טלפון (פורמטים נפוצים)
            phone_patterns = [
                r'\+\d{1,3}[-.\s]?\(?\d{1,3}\)?[-.\s]?\d{1,4}[-.\s]?\d{1,4}',  # פורמט בינלאומי
                r'\(?\d{2,3}\)?[-.\s]?\d{3}[-.\s]?\d{4}',  # פורמט מקומי
                r'\b\d{2,3}[-.\s]?\d{3}[-.\s]?\d{4}\b'  # פורמט פשוט
            ]
            
            phones = []
            for pattern in phone_patterns:
                phones.extend(re.findall(pattern, text_content))
            
            if phones:
                structured_data["phone_numbers"] = list(set(phones))
            
            # חילוץ תאריכים
            date_patterns = [
                r'\d{1,2}/\d{1,2}/\d{2,4}',  # DD/MM/YYYY
                r'\d{1,2}-\d{1,2}-\d{2,4}',  # DD-MM-YYYY
                r'\d{1,2}\.\d{1,2}\.\d{2,4}'  # DD.MM.YYYY
            ]
            
            dates = []
            for pattern in date_patterns:
                dates.extend(re.findall(pattern, text_content))
            
            if dates:
                structured_data["dates"] = list(set(dates))
            
            return structured_data
            
        except Exception as e:
            logger.error(f"שגיאה בחילוץ מידע מובנה מקובץ {file_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def batch_analyze_directory(self, directory_path: str, 
                                recursive: bool = True, 
                                file_pattern: str = "*") -> Dict[str, Any]:
        """
        ניתוח אצווה של תיקייה שלמה
        
        Args:
            directory_path: נתיב לתיקייה
            recursive: האם לבצע ניתוח רקורסיבי
            file_pattern: תבנית שמות קבצים לניתוח (*.txt, *.pdf, וכו')
            
        Returns:
            מילון עם תוצאות הניתוח
        """
        if not os.path.isdir(directory_path):
            return {"status": "error", "error": f"התיקייה {directory_path} אינה קיימת"}
        
        try:
            # איסוף כל הקבצים המתאימים
            import glob
            
            if recursive:
                pattern = os.path.join(directory_path, "**", file_pattern)
                files = glob.glob(pattern, recursive=True)
            else:
                pattern = os.path.join(directory_path, file_pattern)
                files = glob.glob(pattern)
            
            # ניתוח קבצים במקביל
            results = self.analyze_documents(files)
            
            # סיכום התוצאות
            summary = {
                "status": "success",
                "directory": directory_path,
                "file_count": len(files),
                "analyzed_count": len(results),
                "file_types": {},
                "average_file_size": 0,
                "word_count_total": 0,
                "document_analysis": results
            }
            
            # חישוב סטטיסטיקה
            if results:
                # ספירת סוגי קבצים
                file_types = collections.Counter()
                file_sizes = []
                word_count = 0
                
                for result in results:
                    if result.get("status") == "success":
                        mime_type = result.get("mime_type", "unknown")
                        file_types[mime_type] += 1
                        
                        file_size = result.get("file_size", 0)
                        if file_size:
                            file_sizes.append(file_size)
                        
                        word_count += result.get("word_count", 0)
                
                summary["file_types"] = dict(file_types)
                summary["average_file_size"] = sum(file_sizes) / len(file_sizes) if file_sizes else 0
                summary["word_count_total"] = word_count
            
            return summary
            
        except Exception as e:
            logger.error(f"שגיאה בניתוח אצווה של תיקייה {directory_path}: {str(e)}")
            return {"status": "error", "error": str(e)}

    def create_document_topic_model(self, file_paths: List[str], num_topics: int = 5) -> Dict[str, Any]:
        """
        יצירת מודל נושאים ממסמכים
        
        Args:
            file_paths: רשימת נתיבים לקבצים
            num_topics: מספר נושאים לחילוץ
            
        Returns:
            מילון עם תוצאות ניתוח הנושאים
        """
        if not NLP_AVAILABLE:
            logger.warning("ניתוח NLP אינו זמין, לא ניתן ליצור מודל נושאים")
            return {"status": "error", "error": "ניתוח NLP אינו זמין"}
        
        try:
            # חילוץ טקסט מכל המסמכים
            documents = []
            for file_path in file_paths:
                text = self.extract_text_from_document(file_path)
                if text:
                    documents.append({
                        "file_path": file_path,
                        "file_name": os.path.basename(file_path),
                        "text_content": text
                    })
            
            if not documents:
                return {"status": "error", "error": "לא נמצא טקסט במסמכים"}
            
            # יצירת מודל נושאים בסיסי
            from sklearn.decomposition import LatentDirichletAllocation, NMF
            from sklearn.feature_extraction.text import CountVectorizer
            
            # טוקניזציה וסינון מילות עצירה
            vectorizer = CountVectorizer(max_df=0.95, min_df=2, stop_words=self.stop_words)
            
            # התאמת וקטוריזציה למסמכים
            doc_term_matrix = vectorizer.fit_transform([doc["text_content"] for doc in documents])
            
            # הרצת מודל LDA
            lda = LatentDirichletAllocation(
                n_components=num_topics,
                max_iter=10,
                learning_method='online',
                random_state=42
            )
            
            lda.fit(doc_term_matrix)
            
            # שמירת המילים המייצגות כל נושא
            feature_names = vectorizer.get_feature_names_out()
            
            topics = []
            for topic_idx, topic in enumerate(lda.components_):
                top_words_idx = topic.argsort()[:-11:-1]  # 10 המילים המובילות
                top_words = [feature_names[i] for i in top_words_idx]
                topics.append({
                    "id": topic_idx,
                    "top_words": top_words,
                    "weight": float(topic.sum())
                })
            
            # מציאת הנושא הדומיננטי בכל מסמך
            document_topics = lda.transform(doc_term_matrix)
            
            for i, doc in enumerate(documents):
                # נושא דומיננטי
                dominant_topic = int(document_topics[i].argmax())
                doc["dominant_topic"] = dominant_topic
                
                # התפלגות נושאים
                topic_distribution = {}
                for topic_idx, weight in enumerate(document_topics[i]):
                    topic_distribution[f"topic_{topic_idx}"] = float(weight)
                
                doc["topic_distribution"] = topic_distribution
            
            # חישוב דמיון בין מסמכים על בסיס נושאים
            from sklearn.metrics.pairwise import cosine_similarity
            similarity_matrix = cosine_similarity(document_topics)
            
            # אשכול מסמכים דומים
            document_clusters = {}
            for i in range(num_topics):
                # מסמכים עם נושא דומיננטי זהה
                cluster_docs = []
                for j, doc in enumerate(documents):
                    if doc["dominant_topic"] == i:
                        doc_info = {
                            "file_path": doc["file_path"],
                            "file_name": doc["file_name"],
                            "topic_weight": float(document_topics[j][i])
                        }
                        cluster_docs.append(doc_info)
                
                document_clusters[f"topic_{i}"] = cluster_docs
            
            # תוצאות
            result = {
                "status": "success",
                "topics": topics,
                "document_count": len(documents),
                "document_clusters": document_clusters,
                "documents": documents
            }
            
            return result
            
        except Exception as e:
            logger.error(f"שגיאה ביצירת מודל נושאים: {str(e)}")
            return {"status": "error", "error": str(e)}

# פונקציה לשימוש ישיר כמודול עצמאי
def main():
    """פונקציית הרצה עצמאית"""
    import argparse
    
    parser = argparse.ArgumentParser(description='מנתח מסמכים חכם')
    parser.add_argument('file', help='קובץ או תיקייה לניתוח')
    parser.add_argument('--output', default='document_analysis.json', help='קובץ פלט')
    parser.add_argument('--summary', action='store_true', help='יצירת סיכום מסמך')
    parser.add_argument('--extract-text', action='store_true', help='חילוץ טקסט בלבד')
    parser.add_argument('--keywords', action='store_true', help='חילוץ מילות מפתח')
    parser.add_argument('--recursive', action='store_true', help='ניתוח רקורסיבי של תיקייה')
    
    args = parser.parse_args()
    
    # הגדרת לוגים
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )
    
    # יצירת מנתח מסמכים
    analyzer = DocumentAnalyzer()
    
    # בדיקה אם זה קובץ או תיקייה
    if os.path.isdir(args.file):
        print(f"מנתח תיקייה: {args.file}")
        
        if args.recursive:
            print("מבצע ניתוח רקורסיבי...")
        
        # ניתוח תיקייה
        results = analyzer.batch_analyze_directory(args.file, recursive=args.recursive)
        
        print(f"ניתחתי {results.get('analyzed_count', 0)} קבצים מתוך {results.get('file_count', 0)}")
        
    else:
        print(f"מנתח קובץ: {args.file}")
        
        if args.extract_text:
            # חילוץ טקסט בלבד
            text = analyzer.extract_text_from_document(args.file)
            with open(args.output, 'w', encoding='utf-8') as f:
                f.write(text)
            
            print(f"טקסט חולץ ונשמר ב-{args.output}")
            return 0
        
        if args.summary:
            # יצירת סיכום מסמך
            results = analyzer.create_document_summary(args.file)
        elif args.keywords:
            # ניתוח מסמך עם דגש על מילות מפתח
            results = analyzer.analyze_document(args.file)
            
            if "keywords" in results:
                print("מילות מפתח:")
                for keyword in results["keywords"]:
                    print(f"- {keyword}")
        else:
            # ניתוח מסמך רגיל
            results = analyzer.analyze_document(args.file)
    
    # שמירת תוצאות
    with open(args.output, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    
    print(f"התוצאות נשמרו בקובץ: {args.output}")
    return 0

if __name__ == "__main__":
    import time
    sys.exit(main())
